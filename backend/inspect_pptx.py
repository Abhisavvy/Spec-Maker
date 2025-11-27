from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        print(f"Slide {i}: '{title}'")
else:
    print("Template not found")
