from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"--- Slide 0 Analysis ---")
    slide = prs.slides[0]
    if slide.shapes.title:
        print(f"Title Shape Text: '{slide.shapes.title.text}'")
        if slide.shapes.title.has_text_frame:
             for p_idx, paragraph in enumerate(slide.shapes.title.text_frame.paragraphs):
                print(f"  Para {p_idx}: '{paragraph.text}'")
                for r_idx, run in enumerate(paragraph.runs):
                    print(f"    Run {r_idx}: '{run.text}'")
    
    for i, shape in enumerate(slide.shapes):
        print(f"\nShape {i} (Type: {shape.shape_type}):")
        if shape.has_text_frame:
            print(f"  Text: '{shape.text_frame.text}'")
else:
    print("Template not found")
