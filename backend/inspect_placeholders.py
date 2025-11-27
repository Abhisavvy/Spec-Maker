from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"--- Template Analysis ({len(prs.slides)} slides) ---")
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "NO_TITLE"
        print(f"\nSlide {i}: {title}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    print(f"  - Text Shape: '{text}'")
else:
    print("Template not found")
