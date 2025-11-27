from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"--- UI Master Analysis (Slide 8) ---")
    slide = prs.slides[8]
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"\nShape {i}:")
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                print(f"  Para {p_idx}: '{paragraph.text}'")
                for r_idx, run in enumerate(paragraph.runs):
                    print(f"    Run {r_idx}: '{run.text}'")
else:
    print("Template not found")
