from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        # Title
        if slide.shapes.title:
            print(f"TITLE: {slide.shapes.title.text}")
        
        # Text Content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                print(f"CONTENT: {shape.text}")
else:
    print("Template not found")
