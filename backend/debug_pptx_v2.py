from pptx import Presentation
from app.services.pptx_generator import PPTXGenerator
import os

# Mock data simulating LLM output
mock_slides_data = {
    "title": "Debug Spec V2",
    "slides": [
        {
            "header": "Login Screen UI",
            "content": [
                "**Header**: Login",
                "**Sub text**: Please enter your credentials",
                "**CTA**: Login Button",
                "**CTA functionality**: Validates user and redirects to dashboard",
                "**Surfacing conditions**: User is not logged in",
                "**Popup Priority**: High",
                "**Mockup**: https://via.placeholder.com/300"
            ]
        },
        {
            "header": "Onboarding Flow",
            "content": [
                "Step 1: User opens app",
                "**Mockup**: https://via.placeholder.com/150",
                "Step 2: User sees welcome screen",
                "**Mockup**: https://via.placeholder.com/150"
            ]
        }
    ]
}

def debug_pptx():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(base_dir, "../data/slides/Spec Template.pptx")
    output_path = "debug_output_v2.pptx"
    
    if not os.path.exists(template_path):
        print(f"ERROR: Template not found at {template_path}")
        return

    print(f"Using template: {template_path}")
    
    generator = PPTXGenerator()
    
    # Simulate Markdown input
    markdown_sim = """
# Login Screen UI
- **Header**: Login
- **Sub text**: Please enter your credentials
- **CTA**: Login Button
- **CTA functionality**: Validates user and redirects to dashboard
- **Surfacing conditions**: User is not logged in
- **Popup Priority**: High
- **Mockup**: https://via.placeholder.com/300

# Onboarding Flow
- **Description**: Step 1: User opens app. Step 2: User sees welcome screen.
- **Mockup**: https://via.placeholder.com/150

# Vision and anti-vision
- **Vision**: Vision point 1
- **Anti-vision**: Anti-vision point 1
    """
    
    try:
        # Correct signature: create_presentation(markdown_content, output_path, template_path)
        out = generator.create_presentation(markdown_sim, output_path, template_path)
        print(f"Generated: {out}")
        
        # Inspect output
        if os.path.exists(output_path):
            prs = Presentation(output_path)
            print(f"Output has {len(prs.slides)} slides")
            
            # Check for our generated slides
            found_ui = False
            found_flow = False
            
            for i, slide in enumerate(prs.slides):
                title = slide.shapes.title.text if slide.shapes.title else "NO_TITLE"
                print(f"Slide {i}: {title}")
                
                if "Login Screen UI" in title:
                    found_ui = True
                    # Check content
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            txt = shape.text_frame.text
                            if "Login" in txt and "Header" not in txt: # Value filled?
                                print(f"  -> Found filled value: {txt}")
                
                if "Onboarding Flow" in title:
                    found_flow = True

            if found_ui and found_flow:
                print("SUCCESS: Found both generated slides.")
            else:
                print("FAILURE: Missing generated slides.")
        else:
            print("FAILURE: Output file not created.")

    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    debug_pptx()
