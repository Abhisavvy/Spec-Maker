from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    try:
        prs = Presentation(path)
        print(f"Total Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else "NO_TITLE"
            print(f"Slide {i}: {title}")
            # Check for specific markers
            for shape in slide.shapes:
                if hasattr(shape, "text") and "header" in shape.text.lower():
                    print(f"  -> Found 'header' marker in shape: {shape.text[:50]}...")
    except Exception as e:
        print(f"Error reading template: {e}")
else:
    print("Template not found")
