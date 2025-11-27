from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    # Slide indices are 0-based. Slide 9 is index 8, Slide 10 is index 9.
    indices = [8, 9] 
    
    for idx in indices:
        if idx < len(prs.slides):
            slide = prs.slides[idx]
            print(f"\n--- Slide {idx+1} (Layout: {slide.slide_layout.name}) ---")
            if slide.shapes.title:
                print(f"TITLE: '{slide.shapes.title.text}'")
            
            print("SHAPES:")
            for i, shape in enumerate(slide.shapes):
                text = ""
                if hasattr(shape, "text"):
                    text = shape.text
                elif shape.has_text_frame:
                    text = shape.text_frame.text
                
                print(f"  Shape {i}: Type={shape.shape_type}, Name='{shape.name}', Text='{text}'")
else:
    print("Template not found")
