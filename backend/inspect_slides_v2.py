from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"--- Template Analysis ---")
    
    # Vision Slide (Index 3 based on previous run)
    print("\n--- Slide 3 (Vision) ---")
    slide = prs.slides[3]
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Shape {i}: '{shape.text_frame.text}'")
            
    # UI Master (Index 8)
    print("\n--- Slide 8 (UI Master) ---")
    slide = prs.slides[8]
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Shape {i}: '{shape.text_frame.text}'")

    # Flow Master (Index 9)
    print("\n--- Slide 9 (Flow Master) ---")
    slide = prs.slides[9]
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"Shape {i}: '{shape.text_frame.text}'")
else:
    print("Template not found")
