import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import re
import os
from typing import List

class PPTXGenerator:
    def __init__(self):
        pass

    def create_presentation(self, markdown_content: str, output_path: str, template_path: str = None):
        """
        Converts structured Markdown into a PPTX file.
        If template_path is provided, it fills existing slides based on title matching.
        """
        print(f"DEBUG: create_presentation called with template_path={template_path}")
        if template_path and os.path.exists(template_path):
            print("DEBUG: Template found, loading...")
            prs = Presentation(template_path)
            print(f"DEBUG: Loaded template with {len(prs.slides)} slides")
            # Map slide titles to slide objects for easy lookup
            slide_map = {}
            for slide in prs.slides:
                if slide.shapes.title and slide.shapes.title.text:
                    # Normalize title for matching
                    key = slide.shapes.title.text.strip().lower()
                    slide_map[key] = slide
        else:
            print("DEBUG: Template NOT found or not provided, creating new.")
            prs = Presentation()
            slide_map = {}

        # 1. Parse Markdown
        slides_data = self._parse_markdown(markdown_content)
        print(f"DEBUG: Parsed {len(slides_data.get('slides', []))} slides from markdown")
        
        # 3. Pre-fetch Images in Parallel
        # Extract all image URLs from slides_data to download them concurrently
        image_urls = set()
        for slide in slides_data["slides"]:
            for line in slide["content"]:
                # Regex to find markdown images ![alt](url)
                matches = re.findall(r'!\[.*?\]\((.*?)\)', line)
                image_urls.update(matches)
        
        print(f"DEBUG: Found {len(image_urls)} unique images to download.")
        
        image_cache = {}
        if image_urls:
            import concurrent.futures
            
            def download_image(url):
                try:
                    # print(f"DEBUG: Downloading {url}...")
                    resp = requests.get(url, timeout=10)
                    if resp.status_code == 200:
                        return url, resp.content
                except Exception as e:
                    print(f"Error downloading {url}: {e}")
                return url, None

            # Download in parallel
            with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
                future_to_url = {executor.submit(download_image, url): url for url in image_urls}
                for future in concurrent.futures.as_completed(future_to_url):
                    url, content = future.result()
                    if content:
                        image_cache[url] = content
            
            print(f"DEBUG: Successfully pre-fetched {len(image_cache)} images.")

        # 4. Identify Master Slides Dynamically
        ui_master_slide = None
        flow_master_slide = None
        
        for i, slide in enumerate(prs.slides):
            if not slide.shapes.title or not slide.shapes.title.text:
                continue
            title = slide.shapes.title.text.lower().strip()
            
            if "<screen name>" in title or " ui" in title:
                ui_master_slide = slide
                print(f"DEBUG: Found UI Master at index {i}")
            elif "<flow name>" in title or " flow" in title:
                flow_master_slide = slide
                print(f"DEBUG: Found Flow Master at index {i}")

        # 5. Define Template Mapping (LLM Header -> Slide Index)
        # Based on "inspect_titles.py" output
        TEMPLATE_SLIDE_MAPPING = {
            "problem statements": 2,
            "vision": 3,
            "vision and anti-vision": 3,
            "anti-vision": 3,
            "business goals": 4,
            "design goals": 4,
            "opportunity identified": 5,
            "opportunities": 5,
            "expected upsides": 6,
            "overview": 7,
            "edge cases": 10,
            "ui dev requirement": 11,
            "sound requirement": 12,
            "experimentation plan": 13,
            "tracking requirement": 14,
            "analysis plan": 15
        }

        # 6. Fill Slides
        for slide_content in slides_data["slides"]:
            header = slide_content["header"]
            norm_header = header.strip().lower()
            print(f"DEBUG: Processing slide '{header}'")
            
            target_slide = None
            
            # A. Check Explicit Mapping
            for key, idx in TEMPLATE_SLIDE_MAPPING.items():
                if key in norm_header:
                    if idx < len(prs.slides):
                        target_slide = prs.slides[idx]
                        print(f"DEBUG: Mapped '{header}' to Template Slide {idx}")
                        break
            
            # B. Check Master Slides (UI/Flow)
            if not target_slide:
                if (" ui" in norm_header or "screen" in norm_header) and ui_master_slide:
                    print("DEBUG: Cloning UI Master")
                    target_slide = self._duplicate_slide(prs, ui_master_slide)
                    if target_slide.shapes.title:
                        target_slide.shapes.title.text = header
                elif "flow" in norm_header and flow_master_slide:
                    print("DEBUG: Cloning Flow Master")
                    target_slide = self._duplicate_slide(prs, flow_master_slide)
                    if target_slide.shapes.title:
                        target_slide.shapes.title.text = header
            
            # C. Fallback: Fuzzy Match Existing Slides
            if not target_slide:
                # Try to find by title in existing slides
                for i, slide in enumerate(prs.slides):
                    if slide == ui_master_slide or slide == flow_master_slide:
                        continue
                    
                    if slide.shapes.title and slide.shapes.title.text:
                        t_text = slide.shapes.title.text.lower()
                        if norm_header in t_text or t_text in norm_header:
                            target_slide = slide
                            print(f"DEBUG: Keyword matched existing slide {i} for '{header}'")
                            break
            
            # D. Create New Generic Slide
            if not target_slide:
                print(f"DEBUG: Creating NEW slide for '{header}'")
                layout_idx = 1 # Title and Content
                if len(prs.slide_layouts) > 1: layout_idx = 1
                layout = prs.slide_layouts[layout_idx]
                target_slide = prs.slides.add_slide(layout)
                if target_slide.shapes.title:
                    target_slide.shapes.title.text = header
            
            # Fill Content (Pass image_cache)
            self._fill_slide(target_slide, slide_content["content"], image_cache)

        # 7. Save
        prs.save(output_path)
        return output_path

    def _fill_slide(self, slide, content_lines, image_cache=None):
        """Populate a slide with content lines, attempting to fill specific placeholders."""
        
        # 1. Parse content into key-value pairs
        kv_data = {}
        general_content = []
        
        current_key = None
        for line in content_lines:
            # Robust Regex: Matches **Key**: Value, **Key:** Value, - Key: Value, etc.
            match = re.match(r'^\s*[-*]?\s*(.*?)\s*:\s*(.*)', line)
            
            # Check for Sub-header style: **Key** (no colon, or colon at end but empty value)
            sub_header_match = re.match(r'^\s*[-*]?\s*\*\*(.*?)\*\*\s*:?$', line)
            
            if match and match.group(2).strip():
                # Standard Key: Value
                raw_key = match.group(1)
                value = match.group(2).strip()
                key = raw_key.replace("*", "").strip().lower()
                if key.endswith(":"): key = key[:-1].strip()
                kv_data[key] = value
                current_key = key
            elif sub_header_match:
                # Sub-header found (e.g. **Vision**)
                raw_key = sub_header_match.group(1)
                key = raw_key.replace("*", "").strip().lower()
                if key.endswith(":"): key = key[:-1].strip()
                # Start a new key with empty value, expecting following lines to append
                kv_data[key] = ""
                current_key = key
            elif current_key and line.strip():
                # Continuation of previous key
                # Add newline if not empty
                if kv_data[current_key]:
                    kv_data[current_key] += "\n" + line.strip()
                else:
                    kv_data[current_key] = line.strip()
            else:
                general_content.append(line)

        # 2. Iterate shapes and replace
        # Track filled keys to avoid duplicates (especially for Flow slides)
        filled_keys = set()

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text
            lower_text = text.lower()
            
            # --- UI SLIDE SPECIAL HANDLING ---
            # The UI slide has a single large text box with multiple placeholders like "Header : <add header...>"
            # We need to replace these placeholders IN-PLACE.
            if "header :" in lower_text and "sub text :" in lower_text:
                # This is the big UI details box
                # We'll iterate through paragraphs and replace content
                for paragraph in shape.text_frame.paragraphs:
                    p_text = paragraph.text.lower()
                    
                    # Define mappings for this specific text box
                    # Key in KV Data -> Trigger in Template
                    ui_mappings = {
                        "header": "header :",
                        "sub text": "sub text :",
                        "cta": "cta :",
                        "cta functionality": "<add cta functionality here>", # Special case
                        "surfacing conditions": "surfacing conditions :",
                        "popup priority": "popup priority :"
                    }
                    
                    for key, trigger in ui_mappings.items():
                        if trigger in p_text and key in kv_data:
                            val = kv_data[key]
                            
                            # ROBUST REPLACEMENT STRATEGY:
                            # 1. Clear the paragraph.
                            # 2. Re-add the label (if it's a Key: Value type).
                            # 3. Add the value.
                            
                            # Preserve bullet level if any (though usually level 0 for these)
                            level = paragraph.level
                            
                            paragraph.clear()
                            
                            if key == "cta functionality":
                                # This one is just a placeholder line, no label
                                run = paragraph.add_run()
                                run.text = val
                            elif key == "surfacing conditions":
                                # This one has a label "Surfacing conditions :"
                                # But the placeholder might be in the NEXT paragraph?
                                # The template analysis showed:
                                # Para 6: 'Surfacing conditions : '
                                # Para 7: '<add surfacing conditions in bullet points here>'
                                # So if we match Para 6, we just set the label.
                                # But wait, if the user provided value is "Cond 1, Cond 2", we might want to put it here.
                                # Let's just put it on the same line for simplicity and robustness.
                                run_label = paragraph.add_run()
                                run_label.text = "Surfacing conditions : "
                                run_label.font.bold = True
                                
                                run_val = paragraph.add_run()
                                run_val.text = val
                            else:
                                # Standard Key: Value
                                # e.g. "Header : Value"
                                # We want "Header :" to be bold? The template has it plain usually, but let's keep it clean.
                                # Actually, let's try to preserve the style of the first run if possible?
                                # Too complex. Let's just write it out.
                                
                                # Re-construct the label from the trigger (capitalized)
                                label = trigger.capitalize() 
                                if "cta" in label.lower(): label = "CTA :"
                                if "sub text" in label.lower(): label = "Sub text :"
                                if "popup" in label.lower(): label = "Popup Priority :"
                                
                                run_label = paragraph.add_run()
                                run_label.text = label + " "
                                # run_label.font.bold = True # Optional: make labels bold for better look
                                
                                run_val = paragraph.add_run()
                                run_val.text = val
                            
                            paragraph.level = level
                            break # Done with this paragraph
                    
                    # Special handling for the "Surfacing conditions" placeholder paragraph
                    # If we didn't match a key in this paragraph, check if it's the generic placeholder
                    if "<add surfacing conditions" in p_text:
                        # This is the placeholder following the label.
                        # If we already filled it in the previous paragraph (the label one), we should clear this.
                        # OR, if we want to support multi-line bullets here.
                        if "surfacing conditions" in kv_data:
                            paragraph.clear() # Clear it because we appended to the label line above
                            # OR we could put the content here.
                            # Let's clear it to avoid duplication.
                continue # Done with this shape

            # --- STANDARD REPLACEMENTS ---
            # Map specific keys to placeholders found in the template
            replacements = {
                "vision": ["<add vision"],
                "anti-vision": ["<add anti vision"],
                "business goals": ["<add business goals"],
                "design goals": ["<add design goals"],
                "mockup": ["<add mock link here", "mockup"],
                "description": ["<add flow description", "description"]
            }
            
            for key, placeholders in replacements.items():
                if key in kv_data:
                    # Check if this shape matches
                    is_match = False
                    for ph in placeholders:
                        if ph in lower_text:
                            is_match = True
                            break
                    
                    if is_match:
                        val = kv_data[key]
                        
                        # Special handling for Mockups
                        if key == "mockup":
                            if "mockup" not in filled_keys:
                                # Extract URL
                                url = val
                                if "![" in url:
                                    try: url = url.split("](")[1][:-1]
                                    except: pass
                                
                                if url.startswith("http"):
                                    self._replace_shape_with_image(slide, shape, url, image_cache)
                                    filled_keys.add("mockup")
                            else:
                                # Clear duplicate mockup placeholders
                                shape.text_frame.text = ""
                        
                        # Special handling for Flow Descriptions (allow multiple)
                        elif key == "description":
                             if "description" not in filled_keys:
                                 shape.text_frame.text = val
                                 filled_keys.add("description")
                             else:
                                 # For flow slides, we actually WANT multiple descriptions if we have them.
                                 # But our KV data currently only has one "description" key usually?
                                 # If the user provided multiple descriptions in the markdown, they might be in 'general_content'
                                 # or combined in the 'description' key.
                                 # Let's assume for now we clear duplicates to be safe, 
                                 # UNLESS we implement splitting the description.
                                 shape.text_frame.text = ""
                        
                        # Standard Text
                        else:
                            # For lists (Vision, Goals), split by newlines and add paragraphs
                            shape.text_frame.clear()
                            for line in val.split('\n'):
                                p = shape.text_frame.add_paragraph()
                                p.text = line.strip("- *")
                                p.level = 0
                        
                        break # Move to next shape

        # 3. Handle Flow Slides (Advanced Filling)
        # If we have a flow slide, we might want to distribute content across the multiple boxes
        if slide.shapes.title and "flow" in slide.shapes.title.text.lower():
             # Find all description and mockup boxes
             desc_boxes = []
             mock_boxes = []
             for shape in slide.shapes:
                 if shape.has_text_frame:
                     t = shape.text_frame.text.lower()
                     if "description" in t:
                         desc_boxes.append(shape)
                     elif "mock" in t:
                         mock_boxes.append(shape)
            
             # Sort by position (top to bottom, left to right)
             desc_boxes.sort(key=lambda s: (s.top, s.left))
             mock_boxes.sort(key=lambda s: (s.top, s.left))
             
             # If we have general content lines that look like steps, fill them
             step_idx = 0
             
             # Also check if 'description' key has newlines
             all_steps = []
             if "description" in kv_data:
                 all_steps.extend(kv_data["description"].split('\n'))
             all_steps.extend(general_content)
             
             for line in all_steps:
                 if not line.strip(): continue
                 
                 # Check for image in line
                 img_match = re.search(r'!\[.*?\]\((.*?)\)', line)
                 if img_match:
                     # It's an image line
                     if step_idx < len(mock_boxes):
                         self._replace_shape_with_image(slide, mock_boxes[step_idx], img_match.group(1), image_cache)
                         # Don't increment step_idx yet, usually image goes with text?
                         # Actually, in the template, they are pairs.
                         # Let's assume 1 text + 1 image = 1 step
                 else:
                     # It's text
                     if step_idx < len(desc_boxes):
                         desc_boxes[step_idx].text_frame.text = line.strip("- *")
                         # If we successfully filled a text box, we can assume we might fill the corresponding image next
                         # But we need to handle the case where image comes AFTER text
                         pass
            
             # Clean up unused boxes
             # (Optional, but good for polish)

    def _duplicate_slide(self, prs, source_slide):
        """
        Duplicate a slide in the presentation.
        This is tricky in python-pptx as it doesn't have a native clone method.
        We have to create a new slide with the same layout and copy shapes.
        """
        slide_layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Copy shapes
        for shape in source_slide.shapes:
            # We can't easily deep copy shapes with all properties (style, position, etc.)
            # But for placeholders, using the same layout helps.
            # For added textboxes, we need to recreate them.
            
            if shape.is_placeholder:
                # Placeholders are already there due to layout, we just need to copy content if any?
                # Actually, we want the *structure* of the master slide.
                # If the master slide has text in placeholders (like instructions), we might want to copy it
                # so our finder logic works.
                for new_shape in new_slide.placeholders:
                    if new_shape.placeholder_format.idx == shape.placeholder_format.idx:
                        if shape.has_text_frame:
                            new_shape.text = shape.text
            else:
                # Recreate textboxes/shapes
                # This is complex. For now, let's assume the template relies heavily on Layout Placeholders.
                # If there are custom textboxes (like in the UI slide), we need to add them.
                
                if shape.shape_type == 17: # MSO_SHAPE_TYPE.TEXT_BOX
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text = shape.text
                    # Copy basic style if possible (font size etc? too complex for now)
                elif shape.shape_type == 1: # AUTO_SHAPE
                    new_shape = new_slide.shapes.add_shape(
                        shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
                    )
                    if shape.has_text_frame:
                        new_shape.text = shape.text
        
        return new_slide

    def _replace_shape_with_image(self, slide, replace_shape, url, image_cache=None):
        """Download and add image to slide, replacing the placeholder shape while preserving aspect ratio."""
        # print(f"DEBUG: Attempting to replace shape with image from {url}")
        try:
            image_content = None
            
            # Try cache first
            if image_cache and url in image_cache:
                image_content = image_cache[url]
                # print("DEBUG: Using cached image.")
            else:
                # Fallback to download
                print(f"DEBUG: Cache miss for {url}, downloading...")
                response = requests.get(url, timeout=10)
                if response.status_code == 200:
                    image_content = response.content
                else:
                    print(f"DEBUG: Failed to download image. Status code: {response.status_code}")
            
            if image_content:
                image_stream = BytesIO(image_content)
                
                if replace_shape:
                    # Calculate new dimensions to fit within placeholder while preserving aspect ratio
                    # We need the image size first. python-pptx doesn't give it easily from stream before adding.
                    # So we use PIL.
                    from PIL import Image
                    with Image.open(image_stream) as img:
                        img_w, img_h = img.size
                    
                    # Reset stream position
                    image_stream.seek(0)
                    
                    # Placeholder dimensions
                    ph_w = replace_shape.width
                    ph_h = replace_shape.height
                    ph_left = replace_shape.left
                    ph_top = replace_shape.top
                    
                    # Calculate ratios
                    w_ratio = ph_w / img_w
                    h_ratio = ph_h / img_h
                    
                    # Use the smaller ratio to fit entirely within the box
                    scale = min(w_ratio, h_ratio)
                    
                    new_w = int(img_w * scale)
                    new_h = int(img_h * scale)
                    
                    # Center the image
                    new_left = ph_left + (ph_w - new_w) // 2
                    new_top = ph_top + (ph_h - new_h) // 2
                    
                    # Add image
                    slide.shapes.add_picture(image_stream, new_left, new_top, new_w, new_h)
                    
                    # Clear the placeholder text (or remove the shape if possible? No, just clear text)
                    if replace_shape.has_text_frame:
                        replace_shape.text_frame.clear()
                    
                    # print("DEBUG: Image added successfully (Aspect Ratio Preserved)")
                else:
                    # Default placement
                    slide.shapes.add_picture(image_stream, Inches(1), Inches(2), height=Inches(4))
        except Exception as e:
            print(f"Error adding image: {e}")

    def _parse_markdown(self, text: str):
        """
        Parses markdown to extract title (H1) and slides (H2).
        Returns: { "title": str, "slides": [ {"header": str, "content": [str]} ] }
        """
        lines = text.split('\n')
        data = {"title": "", "slides": []}
        current_slide = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Check for headers
            is_h1 = line.startswith('# ')
            is_h2 = line.startswith('## ')
            
            if is_h1 or is_h2:
                header_text = line.lstrip('#').strip()
                
                # If it's the very first header
                if not data["title"] and not data["slides"] and not current_slide:
                    # Heuristic: If it looks like a slide title (UI/Flow) or if the user didn't provide a title block
                    # treat it as a slide.
                    # Actually, let's just use it as title BUT ALSO start a slide if it has content following it.
                    data["title"] = header_text
                    
                    # If the header suggests a specific slide type, we should probably treat it as a slide too.
                    # Or simpler: Just always start a slide. The Title slide in PPTX is usually separate anyway.
                    # If we create a slide for the "Title", it might duplicate the Title Slide.
                    # But if the user's markdown starts with "# Login Screen UI", we definitely want a slide for that.
                    
                    # Let's check if it matches our "Master Slide" keywords
                    lower_header = header_text.lower()
                    if " ui" in lower_header or "flow" in lower_header or "screen" in lower_header:
                         current_slide = {"header": header_text, "content": []}
                    else:
                        # It's likely just the document title
                        pass
                else:
                    if current_slide:
                        data["slides"].append(current_slide)
                    current_slide = {"header": header_text, "content": []}
            elif current_slide:
                # Add content to current slide
                current_slide["content"].append(line)
        
        if current_slide:
            data["slides"].append(current_slide)
            
        return data
