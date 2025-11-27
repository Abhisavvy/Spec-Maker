import os
from typing import List, Dict, Any
import PyPDF2
from docx import Document
import pandas as pd

class DocumentParser:
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """Extract text from a PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {e}")
        return text

    @staticmethod
    def parse_docx(file_path: str) -> str:
        """Extract text from a DOCX file."""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {e}")
        return text

    @staticmethod
    def parse_xlsx(file_path: str) -> str:
        """Extract text from an Excel file (for edge cases)."""
        text = ""
        try:
            df = pd.read_excel(file_path)
            # Convert the dataframe to a string representation
            text = df.to_string(index=False)
        except Exception as e:
            print(f"Error parsing XLSX {file_path}: {e}")
        return text

    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """Extract text from a PPTX file."""
        text = ""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {e}")
        return text

    @staticmethod
    def parse_file(file_path: str) -> str:
        """Dispatch to appropriate parser based on extension."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return DocumentParser.parse_pdf(file_path)
        elif ext == '.docx':
            return DocumentParser.parse_docx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return DocumentParser.parse_xlsx(file_path)
        elif ext == '.pptx':
            return DocumentParser.parse_pptx(file_path)
        else:
            return f"Unsupported file format: {ext}"

    @staticmethod
    def load_documents_from_dir(directory: str) -> List[Dict[str, Any]]:
        """Load and parse all supported documents from a directory."""
        documents = []
        if not os.path.exists(directory):
            return documents
            
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path) and not filename.startswith('.'):
                content = DocumentParser.parse_file(file_path)
                documents.append({
                    "filename": filename,
                    "content": content,
                    "path": file_path
                })
        return documents
