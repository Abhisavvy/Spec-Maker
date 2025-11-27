from pptx import Presentation
import os

path = "../data/slides/Spec Template.pptx"
if os.path.exists(path):
    prs = Presentation(path)
    print(f"--- Template Slide Titles ---")
    for i, slide in enumerate(prs.slides):
        title = "<No Title>"
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text
        print(f"Slide {i}: '{title}'")
        
        # Also check for placeholders to confirm identity
        placeholders = []
        for shape in slide.placeholders:
             if shape.has_text_frame:
                 placeholders.append(shape.text_frame.text[:20] + "...")
        print(f"  Placeholders: {placeholders}")

else:
    print("Template not found")
